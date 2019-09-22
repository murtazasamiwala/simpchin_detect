"""Module to check if a document contains Simplified or Traditional Chinese."""
from zhon import cedict
import re
import os
from os.path import abspath, join
import xlrd
import pptx
import csv
from zipfile import ZipFile
import win32com.client as win32
import shutil
base_path = os.path.dirname(abspath('__file__'))
trad = set(list(cedict.traditional))
simp = set(list(cedict.simplified))
both = set([i for i in trad if i in simp])
all_chinese = cedict.all
avail_exts = ['docx', 'doc', 'pptx', 'xls', 'xlsx', 'csv', 'txt', 'rtf']
passed_exts = ['py', 'git', 'spec', 'exe', 'md', 'gitattributes', 'gitignore']


def extract_chinese(fname, path=base_path):
    """Extract Chinese text from given document."""
    exten = fname.split('.')[-1]
    if exten in ['doc', 'docx', 'rtf']:
        word = win32.Dispatch('Word.Application')
        doc_file = path + '\\' + fname
        doc = word.Documents.Open(doc_file)
        txt = doc.Content.Text
        doc.Close(False)
        word.Quit()
    elif (fname.endswith('.xls')) | (fname.endswith('.xlsx')):
        workbook = xlrd.open_workbook(fname)
        sheets_name = workbook.sheet_names()
        txt = '\n'
        for names in sheets_name:
            worksheet = workbook.sheet_by_name(names)
            num_rows = worksheet.nrows
            num_cells = worksheet.ncols
            for curr_row in range(num_rows):
                new_output = []
                for index_col in range(num_cells):
                    value = worksheet.cell_value(curr_row, index_col)
                    if value:
                        new_output.append(value)
                if new_output:
                    txt += ' '.join(new_output) + '\n'
    elif fname.endswith('.pptx'):
        presentation = pptx.Presentation(fname)
        text_runs = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
        txt = '\n\n'.join(text_runs)
    elif fname.endswith('.txt'):
        text_doc = open(fname, 'r', encoding='utf8')
        txt = text_doc.read()
    elif fname.endswith('.csv'):
        csv_doc = open(fname, 'r', encoding='utf8')
        csv_reader = csv.reader(csv_doc, delimiter=',')
        txt = '\n'.join(['\t'.join(row) for row in csv_reader])
    chinese_text = set(re.sub('[^%s]' % all_chinese, '', txt))
    return chinese_text


def market_check(unzip_path):
    """Identify market name based on presence of certain files."""
    market = 'Not known'
    check_name = 'Guidelines_for_identifying_use_of_SC_in_TC_jobs.docx'
    if 'Reference_files' in os.listdir(unzip_path):
        ref_path = unzip_path + '\\' + 'Reference_files'
        for i in os.listdir(ref_path):
            test_name = ''.join([j + '_' for j in i.split('_')[2:]])
            if test_name[:len(test_name)-1] == check_name:
                market = 'Taiwan'
    return market


def zip_extract(zname, path=base_path):
    """Extract files from zip and perform chinese text checks."""
    temp_msg_list = []
    zf = ZipFile(zname)
    orig_dir = os.listdir(path)
    zf.extractall()
    for i in os.listdir(path):
        if i not in orig_dir:
            extension = i.split('.')[-1]
            if extension == i:
                unzip_path = base_path + '\\' + i
                doc_path = unzip_path + '\\' + zname.split('.')[0]
                market_name = market_check(unzip_path)
                for j in os.listdir(doc_path):
                    ext2 = j.split('.')[-1]
                    if ext2 in avail_exts:
                        chinese_text = extract_chinese(doc_path + '\\' + j)
                        msg = chinese_text_check(chinese_text, j, market_name)
                        temp_msg_list.append(msg)
                shutil.rmtree(i)
    zf.close()
    if (market_name == 'Taiwan') & ('ERROR' in msg):
        os.remove(zname)
    return temp_msg_list


def chinese_text_check(ch_text, fname, market):
    """Perform identification checks and return message."""
    if len(ch_text) == 0:
        msg = '{} does not contain Chinese text'.format(fname)
        fin_msg = report(msg, fname, 'IGNORE FILE')
    elif ch_text.issubset(trad):
        if market == 'Taiwan':
            msg_1 = '{} is written in Traditional Chinese'.format(fname)
            msg_2 = ' and market is Taiwan.\nDeliver the job.'
            fin_msg = report(msg_1 + msg_2, fname, 'PASSED')
        else:
            msg_1 = '{} is written in Traditional Chinese.\n'.format(fname)
            msg_2 = 'Confirm that the service is E to TC. '
            msg_3 = 'Otherwise, it is a serious error.'
            comb_msg = msg_1 + msg_2 + msg_3
            fin_msg = report(comb_msg, fname, 'TRADITIONAL CHINESE')
    elif ch_text.issubset(simp):
        if market == 'Taiwan':
            msg_1 = '{} is written in Simplified Chinese'.format(fname)
            msg_2 = ' but market is Taiwan.\n'
            msg_3 = 'This is a serious error!! DO NOT DELIVER!!\n'
            msg_4 = 'File has been deleted!!'
            comb_msg = msg_1 + msg_2 + msg_3 + msg_4
            fin_msg = report(comb_msg, fname, 'ERROR')
        else:
            msg_1 = '{} is written in Simplified Chinese.\n'.format(fname)
            msg_2 = 'Confirm that the service is E to SC. '
            msg_3 = 'Otherwise, it is a serious error.'
            comb_msg = msg_1 + msg_2 + msg_3
            fin_msg = report(comb_msg, fname, 'SIMPLIFIED CHINESE')
    else:
        output_name = 'output_' + fname.split('.')[0] + '.txt'
        output = open(output_name, 'a', encoding='utf8')
        for char in ch_text:
            if char in simp:
                output.write(char + '\n')
        output.close()
        msg_1 = '{} has both Simplified and Traditional '.format(fname)
        msg_2 = 'characters\nCheck service name and fix characters'
        msg_3 = ' of other language.\nFile has been deleted!!\n'
        msg_strt = msg_1 + msg_2 + msg_3
        msg_end_1 = ' has been generated. It is a list of Simplified '
        msg_end_2 = 'Characters to be fixed.'
        msg_end = msg_end_1 + msg_end_2
        full_msg = msg_strt + output_name + msg_end
        fin_msg = report(full_msg, fname, 'ERROR')
    return fin_msg


def report(msg, filename, result):
    """Format for writing to result file."""
    msg_head = '*' * 20 + '\n' + 'Result for {}:'.format(filename) + '\n'
    result_msg = 'RESULT :: '+result+'\n'
    msg_body = msg + '\n' + '-' * 20 + '\n'
    return msg_head + result_msg + msg_body


def directory_check(path=base_path):
    """Run whole module on a given directory."""
    msg_list = []
    market = 'Not known'
    for i in os.listdir(path):
        extension = i.split('.')[-1]
        if extension == i:
            pass
        elif extension in passed_exts:
            pass
        elif extension == 'ppt':
            msg_1 = 'ppt format not supported. If file is in zip, extract it.'
            msg_2 = '\nThen convert {} to pptx and run script.'.format(i)
            msg = msg_1 + msg_2
            msg_list.append(report(msg, i, 'NOT SUPPORTED'))
        elif extension in avail_exts:
            chinese_text = extract_chinese(i)
            msg = chinese_text_check(chinese_text, i, market)
            msg_list.append(msg)
            if 'ERROR' in msg:
                os.remove(i)
        elif extension == 'zip':
            zip_msg = zip_extract(i)
            msg_list = msg_list+zip_msg
        else:
            msg = '{} is not one of the acceptable formats.'.format(i)
            msg_list.append(report(msg, i, 'IGNORE FILE'))
    result = open('script_result.txt', 'a', encoding='utf8')
    for i in msg_list:
        result.write(i)
    result.close()


if __name__ == '__main__':
    directory_check()
